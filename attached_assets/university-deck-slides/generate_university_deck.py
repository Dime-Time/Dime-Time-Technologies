"""Dime Time — University Speaking Deck (founder journey, academic audience).

13 core slides + 1 optional classroom-discussion slide = 14 total. 16:9, native
editable text, embedded speaker notes on every slide. Built per
UNIVERSITY-DECK-BRIEF.txt; trimmed from 21 to 14 slides at founder request
(2026-07-21) and Slide 7 early-foundation milestone left undated per the
final timeline correction.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

PURPLE = RGBColor(0x91, 0x8E, 0xF4)
DEEP = RGBColor(0x4A, 0x46, 0x8F)
DARK = RGBColor(0x2D, 0x2A, 0x4A)
LIGHT = RGBColor(0xF5, 0xF4, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xD6, 0x3A, 0x3A)
GREEN = RGBColor(0x2E, 0x9E, 0x62)
GOLD = RGBColor(0xE8, 0xB8, 0x4B)
LAV2 = RGBColor(0xC9, 0xC7, 0xF9)

HERE = os.path.dirname(__file__)
LOGO = os.path.join(HERE, "..", "patent-deck-slides", "_logo_transparent.png")
SHOT = os.path.join(HERE, "..", "Dime-Time-App-Store-Screenshots-FINAL", "iPhone-6.9")

TOTAL = 14

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height

_n = 0


def slide(bg=LIGHT):
    global _n
    _n += 1
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT, first=False, space_after=6, font="Poppins"):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font
    return p


def chip(s, x, y, w, h, fill, text, size, tcolor=WHITE, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE, sub=None, subsize=12):
    c = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = tcolor; r.font.name = "Poppins"
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(subsize); r2.font.color.rgb = tcolor; r2.font.name = "Poppins"
    return c


def header(s, kicker, title, accent=PURPLE, kcolor=DEEP, tcolor=DARK, tsize=32):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.22))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = box(s, 0.7, 0.42, 12, 0.45)
    para(tf, kicker.upper(), 14, kcolor, bold=True, first=True)
    tf2 = box(s, 0.7, 0.82, 12.0, 0.95)
    para(tf2, title, tsize, tcolor, bold=True, first=True)


def footer(s, light=False):
    tf = box(s, 10.9, 7.06, 2.0, 0.36)
    para(tf, f"{_n:02d} / {TOTAL}", 14, (LAV2 if light else DEEP), bold=True, align=PP_ALIGN.RIGHT, first=True)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ================= SLIDE 1 — COVER =================
s = slide(PURPLE)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.87), Inches(0.5), height=Inches(1.6))
tf = box(s, 0.6, 2.35, 12.1, 0.8)
para(tf, "FROM SOMMELIER TO FINTECH FOUNDER", 32, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
tf = box(s, 0.8, 3.25, 11.7, 0.6)
para(tf, "How a Nontechnical Founder Built and Launched Dime Time", 24, DARK, bold=True, align=PP_ALIGN.CENTER, first=True)
tf = box(s, 1.5, 3.95, 10.3, 0.7)
para(tf, "A real-world journey through entrepreneurship, information systems, financial infrastructure and AI-assisted development", 16, WHITE, align=PP_ALIGN.CENTER, first=True)
tf = box(s, 0.8, 5.0, 11.7, 1.5)
para(tf, "Tim Carlisle", 24, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=2)
para(tf, "Founder & CEO, Dime Time Technologies LLC", 17, WHITE, align=PP_ALIGN.CENTER, space_after=2)
para(tf, "U.S. Navy Veteran  •  MBA • Graduate degree in Management Information Systems", 16, DARK, bold=True, align=PP_ALIGN.CENTER)
tf = box(s, 0.8, 6.6, 11.7, 0.5)
para(tf, "\u201cGet out of debt, one dime at a time.\u201d", 18, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
footer(s, light=True)
notes(s, "Welcome and thank you for having me. My name is Tim Carlisle, founder of Dime Time. Quick show of hands before we start: who here has ever built something — a business, an app, a side project? Keep that feeling in mind, because this talk is about what happens between the idea and the launch.\n\nSet expectations: this is NOT a pitch. It's the honest story of how a Navy veteran and sommelier with no software engineering background shipped a live fintech app — the wins, the rejections, and what I'd tell my younger self sitting where you are.\n\n[20-MIN VERSION] Keep intro under 60 seconds.\n[45-60 MIN] Add a short story about a night in the restaurant that taught you something about pressure.\n[DO NOT OVERSTATE] The app is live; public money movement is not yet enabled — we'll label that carefully later.")

# ================= SLIDE 2 — WHO I AM =================
s = slide(LIGHT)
header(s, "Who I Am", "A Nontraditional Path Into Technology")
items = [
    "U.S. Navy veteran",
    "Hospitality career in Las Vegas and Dallas",
    "Experience opening and operating high-profile restaurants",
    "Sommelier serving demanding, high-expectation guests",
    "MBA • Graduate degree in Management Information Systems",
    "Founder and builder of Dime Time",
]
tf = box(s, 0.9, 2.05, 7.3, 4.4)
for i, t in enumerate(items):
    para(tf, "•  " + t, 20, DARK, first=(i == 0), space_after=12)
pnl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(2.1), Inches(4.1), Inches(4.2))
pnl.fill.solid(); pnl.fill.fore_color.rgb = DEEP
pnl.line.fill.background(); pnl.shadow.inherit = False
tf = box(s, 8.85, 2.5, 3.4, 3.5)
para(tf, "CORE MESSAGE", 14, LAV2, bold=True, first=True, space_after=10)
para(tf, "Technology founders do not all begin as software engineers.", 19, WHITE, bold=True, space_after=10)
para(tf, "Operational experience, customer empathy, discipline and systems thinking can become founder advantages.", 17, LAV2)
footer(s)
notes(s, "Walk through the path briefly: Navy, then nearly two decades in hospitality — Las Vegas and Dallas, opening high-profile restaurants, working as a sommelier. Then the MBA and a graduate degree in Management Information Systems.\n\nThe point to land: none of these jobs had 'software engineer' in the title, and that turned out to be fine. Every stop added a skill the company needed later — discipline, service under pressure, systems thinking.\n\n[PAUSE HERE] Ask: 'How many of you feel like your background is *wrong* for what you want to do?' Let a few hands go up.\n[DO NOT OVERSTATE] Say exactly 'MBA and a graduate degree in Management Information Systems' — do not shorten to 'M.S.'\n[20-MIN VERSION] Compress to 45 seconds — the core message box is the takeaway.")

# ================= SLIDE 3 — CLASSROOM FOUNDATION =================
s = slide(LIGHT)
header(s, "The Classroom Foundation", "The Idea Began as an Academic Exercise")
tf = box(s, 0.9, 2.0, 6.9, 4.5)
para(tf, "Business education turned an observation into a structured opportunity.", 20, DARK, bold=True, first=True, space_after=12)
para(tf, "Entrepreneurship coursework introduced: problem definition, customer value propositions, market analysis, branding, business models, pitch development, financial planning.", 18, DARK, space_after=12)
para(tf, "Management Information Systems education connected business strategy to technology architecture.", 18, DARK, space_after=12)
para(tf, "Early academic work included the Dime Time concept, branding, color scheme, presentation materials and an initial outsourced demo.", 18, DARK)
pnl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(2.1), Inches(4.4), Inches(4.2))
pnl.fill.solid(); pnl.fill.fore_color.rgb = WHITE
pnl.line.color.rgb = PURPLE; pnl.line.width = Pt(2); pnl.shadow.inherit = False
tf = box(s, 8.55, 2.45, 3.7, 3.6)
para(tf, "IMPORTANT FRAMING", 14, DEEP, bold=True, first=True, space_after=10)
para(tf, "The early demo was not the final product.", 19, DARK, bold=True, space_after=10)
para(tf, "It was a learning tool that transformed an abstract idea into something people could see and discuss.", 17, DARK)
footer(s)
notes(s, "This is the slide that connects directly to the students in the room: Dime Time started as coursework. The frameworks they're learning right now — problem definition, value proposition, market analysis — are the exact tools that shaped this company.\n\nBe honest about the early demo: it was outsourced, it was rough, and it was NOT the product that shipped. Its value was making the idea tangible enough to discuss and refine.\n\n[PAUSE HERE] 'Some of you have a project like this sitting in a folder from last semester. What would it take to pick it back up?'\n[DO NOT OVERSTATE] Do not attach a specific year to the classroom concept unless you choose to share it live — no invented dates.")

# ================= SLIDE 4 — THE PROBLEM =================
s = slide(WHITE)
header(s, "The Problem", "Why Dime Time Needed to Exist")
tf = box(s, 0.9, 2.0, 6.7, 4.0)
for i, t in enumerate([
    "Many consumers intend to make extra debt payments",
    "Manual extra payments require repeated discipline",
    "Small amounts are easier to commit than large, irregular payments",
    "Everyday transactions can create a recurring behavioral trigger",
    "Automation can help translate intention into action",
]):
    para(tf, "•  " + t, 19, DARK, first=(i == 0), space_after=13)
stats = [("$18.8T", "U.S. household debt"), ("$1.25T", "credit-card balances"), ("91%", "of U.S. adults own a smartphone")]
y = 1.95
for big, small in stats:
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(y), Inches(4.6), Inches(1.28))
    c.fill.solid(); c.fill.fore_color.rgb = DEEP
    c.line.fill.background(); c.shadow.inherit = False
    tf2 = c.text_frame; tf2.word_wrap = True; tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = big
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Poppins"
    p2 = tf2.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = small
    r2.font.size = Pt(15); r2.font.color.rgb = LAV2; r2.font.name = "Poppins"
    y += 1.48
tf = box(s, 0.9, 6.25, 11.8, 0.75)
para(tf, "These figures describe the scale of the problem — they are not Dime Time metrics.", 16, DEEP, bold=True, first=True, space_after=3)
para(tf, "Sources: Federal Reserve Bank of New York, Household Debt & Credit Report, Q1 2026 (released May 2026)  •  Pew Research Center, Mobile Fact Sheet, 2025", 14, DARK)
footer(s)
notes(s, "Frame the problem behaviorally, not just numerically: almost everyone in debt INTENDS to pay extra. The failure point is that manual extra payments require a fresh act of discipline every single month. Small, automatic amounts remove that friction.\n\nThen give scale: $18.8 trillion in household debt, $1.25 trillion on credit cards alone — NY Fed data from May 2026. And 91% of adults carry the delivery device (a smartphone) in their pocket.\n\n[DO NOT OVERSTATE] Say clearly these are national statistics, NOT Dime Time traction. Students respect the honesty.\n[45-60 MIN] Ask the class: 'Why do you think savings round-ups took off years ago, but debt round-ups didn't?'")

# ================= SLIDE 5 — THE SOLUTION =================
s = slide(LIGHT)
header(s, "The Solution", "Get Out of Debt, One Dime at a Time")
steps = [
    "A user connects an eligible financial account",
    "Everyday purchases are analyzed",
    "Each purchase is rounded up",
    "Round-ups accumulate",
    "The user controls the destination and schedule",
    "Authorized ACH payments apply accumulated amounts toward debt",
]
tf = box(s, 0.9, 2.0, 7.4, 3.9)
for i, t in enumerate(steps):
    para(tf, f"{i+1}.  {t}", 19, DARK, bold=(i == 5), first=(i == 0), space_after=11)
sh1 = os.path.join(SHOT, "01-dashboard-iphone.png")
sh2 = os.path.join(SHOT, "03-roundups-iphone.png")
if os.path.exists(sh1):
    s.shapes.add_picture(sh1, Inches(8.55), Inches(1.95), height=Inches(4.45))
if os.path.exists(sh2):
    s.shapes.add_picture(sh2, Inches(10.75), Inches(1.95), height=Inches(4.45))
pnl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(6.05), Inches(11.6), Inches(0.95))
pnl.fill.solid(); pnl.fill.fore_color.rgb = WHITE
pnl.line.color.rgb = DEEP; pnl.line.width = Pt(1.5); pnl.shadow.inherit = False
tf = box(s, 1.2, 6.2, 11.1, 0.7)
para(tf, "Accuracy note: money does not move after every purchase — round-ups accumulate and transfer on a user-approved schedule.", 16, DEEP, bold=True, first=True)
footer(s)
notes(s, "Walk the six steps slowly — this is the whole product in one slide. The screenshots on the right are the real, live App Store product, not mockups.\n\nEmphasize step 5: the USER controls destination and schedule. This is a consent-first system.\n\nThe accuracy note matters in a university setting: money does not move per-purchase. Round-ups accumulate and transfer on an approved schedule. Getting details like this right is part of operating a financial product responsibly — that's a teaching point in itself.\n\n[20-MIN VERSION] 60–90 seconds; the slogan is the takeaway.\n[PAUSE HERE] Good moment for clarifying questions about how round-ups work.")

# ================= SLIDE 6 — SAFETY BEFORE SCALE =================
s = slide(LIGHT)
header(s, "Safety Before Scale", "Building Financial Technology Means Building Guardrails")
guards = [
    "Master real-transfer switch", "Separate ACH feature switch", "Founder-controlled user allowlist",
    "No automatic public approval", "First-transfer limit", "Daily transfer limit",
    "Transfer-frequency limit", "User authorization evidence", "Duplicate-payment protection",
    "Webhook verification", "Transaction ledger", "Audit logging",
]
x0, y0 = 0.9, 2.0
for i, g in enumerate(guards):
    col = i % 3
    row = i // 3
    chip(s, x0 + col * 3.95, y0 + row * 0.92, 3.75, 0.75, WHITE, g, 14.5, tcolor=DARK, bold=True)
    sh = s.shapes[-1]; sh.line.color.rgb = PURPLE; sh.line.width = Pt(1.5)
tf = box(s, 0.9, 5.8, 11.8, 0.5)
para(tf, "Plus: the ability to disable real transfers instantly.", 16, DARK, bold=True, first=True)
bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(6.3), Inches(11.6), Inches(0.85))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN
bar.line.fill.background(); bar.shadow.inherit = False
tf = box(s, 1.2, 6.47, 11.0, 0.6)
para(tf, "In fintech, the safest successful test is more valuable than an uncontrolled launch.", 17, WHITE, bold=True, first=True)
footer(s)
notes(s, "This is the slide that separates a fintech from a demo. Every card on screen is a real control in the live system.\n\nExplain the philosophy in plain language: real money movement sits behind multiple independent switches. Even with everything on, a user must be individually approved, the first transfer is capped tiny, daily amounts and frequency are limited, and every decision — approve or block — is written to an audit log.\n\nThe payoff line: when we ran the first real test, the question wasn't 'will it work?' but 'which of our safeguards will tell us if it doesn't?'\n\n[DO NOT OVERSTATE] Keep this at the philosophy level — no exact implementation details, code, or credentials.\n[20-MIN VERSION] 60 seconds: 'multiple switches, individual approval, tiny limits, full audit trail.'")

# ================= SLIDE 7 — MILESTONE TIMELINE =================
s = slide(WHITE)
header(s, "The Historical Timeline", "From Academic Concept to Live Fintech Product", tsize=30)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.62), Inches(2.1), Pt(2.5), Inches(4.35))
ln.fill.solid(); ln.fill.fore_color.rgb = LAV2
ln.line.fill.background(); ln.shadow.inherit = False
rows = [
    ("EARLY FOUNDATION", "Classroom concept and initial prototype — entrepreneurship and MIS education shaped the concept, brand and early demo", LAV2, DARK),
    ("Aug 3, 2025", "First commit — Dime Time becomes a software project", DEEP, WHITE),
    ("May 27, 2026", "Privacy Policy & Terms effective — legal foundation prepared for production", DEEP, WHITE),
    ("Jun 29, 2026", "Dime Time goes LIVE on the Apple App Store", PURPLE, WHITE),
    ("Jul 7, 2026", "Controlled $1.00 real ACH debt-payment test initiated through Stripe", PURPLE, WHITE),
    ("Jul 21, 2026", "v1.0.5 (build 207) live  •  $0.99 payout confirmed in Mercury — bank-to-bank loop proven  •  Stripe Financial Connections submitted", GREEN, WHITE),
]
y = 1.98
for date, desc, fill, tc in rows:
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.55), Inches(y + 0.26), Inches(0.18), Inches(0.18))
    d.fill.solid(); d.fill.fore_color.rgb = fill if fill != LAV2 else DEEP
    d.line.fill.background(); d.shadow.inherit = False
    chip(s, 1.95, y, 1.85, 0.62, fill, date, 15 if len(date) <= 12 else 14, tcolor=tc)
    tf = box(s, 4.05, y + 0.02, 8.5, 0.72)
    para(tf, desc, 15, DARK, first=True, space_after=0)
    y += 0.76
tf = box(s, 0.9, 6.68, 11.8, 0.55)
para(tf, "Today: iOS live  •  rails proven in a controlled founder test  •  Google Play prep underway  •  public bank linking & debt import awaiting external approvals", 14, DEEP, bold=True, first=True)
footer(s)
notes(s, "Walk this top to bottom — it's the spine of the whole talk.\n\nOn the first milestone, keep it undated: 'Dime Time's early foundation developed during my academic journey. Entrepreneurship coursework helped me structure the problem, customer value proposition, brand and business model. The original prototype was not the product that exists today — it was a learning tool that made the idea tangible and allowed me to begin testing and communicating the concept.'\n\nThen the dated record: first commit August 3, 2025. Live on the App Store June 29, 2026 — under eleven months later. First real dollar moved July 7. Loop proven bank-to-bank July 21.\n\nThe teaching point: the idea waited for its builder. What changed wasn't the idea; it was the decision to start, plus tools that made a solo build feasible.\n\n[PAUSE HERE] Let the eleven-months-to-launch arc sink in before moving on.\n[DO NOT OVERSTATE] Do not assign a year to the classroom era, do not imply continuous development before the first commit, and do not present the early prototype as the production app. The dated milestones on this slide are the verified ones.")

# ================= SLIDE 8 — ONE DOLLAR =================
s = slide(DARK)
header(s, "The Moment the Model Became Real", "One Dollar Changed the Project", accent=GREEN, kcolor=LAV2, tcolor=WHITE)
steps = [("$1.00", "ACH debt payment"), ("STRIPE", "processing"), ("$0.99", "payout"), ("MERCURY", "business bank"), ("PROVEN", "bank-to-bank loop")]
x = 0.8
bw, gap = 2.15, 0.32
for i, (big, small) in enumerate(steps):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.15), Inches(bw), Inches(1.5))
    c.fill.solid(); c.fill.fore_color.rgb = GREEN if big == "PROVEN" else PURPLE
    c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = big
    r.font.size = Pt(21); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Poppins"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = small
    r2.font.size = Pt(14); r2.font.color.rgb = WHITE; r2.font.name = "Poppins"
    if i < len(steps) - 1:
        tf2 = box(s, x + bw - 0.05, 2.62, gap + 0.12, 0.55)
        para(tf2, "→", 22, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    x += bw + gap
chip(s, 3.97, 4.05, 5.4, 0.7, GREEN, "PROVEN IN CONTROLLED LIVE TEST", 18)
tf = box(s, 1.4, 5.0, 10.5, 2.0)
for i, t in enumerate([
    "The amount was intentionally small — and the test used the founder's own controlled, allowlisted account",
    "The payment moved through real production rails; the payout confirmed the underlying system works",
    "This did NOT mean public transfers were enabled — public bank linking still awaits external approval",
]):
    para(tf, "•  " + t, 16, LAV2 if i < 2 else GOLD, bold=(i == 2), first=(i == 0), space_after=9)
footer(s, light=True)
notes(s, "Slow down here — this is the emotional center of the talk.\n\nTell it as a story: July 7, you pressed the button on a $1.00 payment from your own bank account. ACH takes days and reports back through webhooks, so you waited. It settled. And on July 21, $0.99 — after processing costs — appeared in the company's business account at Mercury, with the company name on a real bank statement.\n\nOne dollar proved the entire machine: consent, debit, processing, settlement, payout, reconciliation.\n\n[DO NOT OVERSTATE] Use the exact phrase 'proven in controlled live test.' It is NOT public money movement — the third bullet is mandatory, not optional.\n[PAUSE HERE] 'Why do you think we tested with one dollar instead of a hundred?' — great discussion moment.\n[20-MIN VERSION] Keep the story to 90 seconds: pressed the button July 7, $0.99 landed at the bank July 21, loop proven — and public transfers remain off.")

# ================= SLIDE 9 — SETBACKS TABLE =================
s = slide(WHITE)
header(s, "Setbacks as System Feedback", "Every Failure Became a Better Control")
cols = [(0.7, 3.3), (4.15, 4.35), (8.65, 4.0)]
heads = ["CHALLENGE", "RESPONSE", "LESSON"]
for (x, w), h in zip(cols, heads):
    chip(s, x, 1.92, w, 0.5, DEEP, h, 15)
table = [
    ("App Store duplicate-app rejection", "Removed the obsolete record and clarified the production app", "Compliance problems require evidence and precise communication"),
    ("Production email failures", "Diagnosed delivery paths, corrected notifications, added tests", "A feature is not working merely because the interface says it is"),
    ("Real-money risk", "Added feature switches, allowlisting, limits and audit controls", "Safety must be designed before activation"),
    ("OAuth return-flow issues", "Added universal-link and resume handling", "The complete user journey matters, not only the API call"),
    ("Android permission risk", "Removed permissions the product did not need", "Minimize access and justify every permission"),
]
y = 2.56
for i, (a, b, c) in enumerate(table):
    if i % 2 == 0:
        band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y - 0.04), Inches(12.2), Inches(0.86))
        band.fill.solid(); band.fill.fore_color.rgb = LIGHT
        band.line.fill.background(); band.shadow.inherit = False
    tf = box(s, cols[0][0] + 0.1, y, cols[0][1] - 0.15, 0.85)
    para(tf, a, 14, DARK, bold=True, first=True, space_after=0)
    tf = box(s, cols[1][0] + 0.1, y, cols[1][1] - 0.15, 0.85)
    para(tf, b, 14, DARK, first=True, space_after=0)
    tf = box(s, cols[2][0] + 0.1, y, cols[2][1] - 0.15, 0.85)
    para(tf, c, 14, DEEP, bold=True, first=True, space_after=0)
    y += 0.86
tf = box(s, 0.7, 6.95, 11.8, 0.4)
para(tf, "Presented as learning and responsible iteration — not blame.", 14, DEEP, first=True)
footer(s)
notes(s, "This is the most honest slide in the deck, and usually the one professors appreciate most.\n\nWalk one or two rows fully. The email failure row lands well: the interface said emails were sent; users weren't receiving them. Lesson — 'the interface says so' is not evidence. That principle now shapes how everything gets verified.\n\nThe pattern to draw out: each failure produced a CONTROL, not just a fix. Rejection → better compliance process. Money risk → layered safeguards. Permissions → minimalism by default.\n\n[PAUSE HERE] Ask students: 'Which of these five failures would you have predicted in advance?' Most predict none — that's the point.\n[20-MIN VERSION] One row only, 60 seconds.")

# ================= SLIDE 10 — AI ROLE =================
s = slide(LIGHT)
header(s, "What AI Did and Did Not Do", "AI as a Force Multiplier, Not an Autopilot")
c1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.95), Inches(5.7), Inches(4.05))
c1.fill.solid(); c1.fill.fore_color.rgb = WHITE
c1.line.color.rgb = PURPLE; c1.line.width = Pt(2); c1.shadow.inherit = False
tf = box(s, 1.25, 2.15, 5.0, 3.7)
para(tf, "WHAT AI HELPED WITH", 15, DEEP, bold=True, first=True, space_after=8)
for t in ["Explaining technical concepts", "Drafting implementation plans", "Reviewing code and configuration", "Generating tests", "Diagnosing failures", "Preparing documentation", "Structuring compliance checklists", "Accelerating iteration"]:
    para(tf, "•  " + t, 15.5, DARK, space_after=5)
c2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.95), Inches(5.7), Inches(4.05))
c2.fill.solid(); c2.fill.fore_color.rgb = DEEP
c2.line.fill.background(); c2.shadow.inherit = False
tf = box(s, 7.25, 2.15, 5.0, 3.7)
para(tf, "WHAT STILL REQUIRED THE FOUNDER", 15, GOLD, bold=True, first=True, space_after=8)
for t in ["Choosing the problem", "Defining the product", "Making risk decisions", "Opening provider accounts", "Identity and business verification", "Reviewing financial behavior", "Running controlled live tests", "Deciding when NOT to launch", "Accepting responsibility for the outcome"]:
    para(tf, "•  " + t, 15.5, WHITE, space_after=4)
bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(6.25), Inches(11.6), Inches(0.9))
bar.fill.solid(); bar.fill.fore_color.rgb = PURPLE
bar.line.fill.background(); bar.shadow.inherit = False
tf = box(s, 1.2, 6.4, 11.0, 0.65)
para(tf, "The advantage is not access to AI — it is the ability to direct it, question it, verify it and combine it with real-world judgment.", 16, WHITE, bold=True, first=True)
footer(s)
notes(s, "Students will lean in here — many are wondering whether AI makes their skills obsolete. Give them the honest split on screen.\n\nLeft column: AI genuinely accelerated everything — explanation, planning, code review, diagnosis. It compressed years of learning into months.\n\nRight column: every consequential decision stayed human. No AI can open your Stripe account, pass identity verification, decide acceptable risk for real money, or take responsibility when something breaks. 'Deciding when NOT to launch' is worth emphasizing — restraint was a founder decision, repeatedly.\n\n[DO NOT OVERSTATE] Never say AI built the company. 'AI-assisted, founder-directed' is the accurate frame.\n[PAUSE HERE] Expect and welcome questions — this generates the most Q&A of any slide.")

# ================= SLIDE 11 — LESSONS =================
s = slide(DARK)
header(s, "Lessons for Students", "What I Wish More Students Knew", accent=PURPLE, kcolor=LAV2, tcolor=WHITE)
lessons = [
    "You do not need permission to begin learning.",
    "Your previous career may contain your founder advantage.",
    "A prototype starts a conversation; a production system earns trust.",
    "Learn enough technology to ask better questions.",
    "Evidence is stronger than confidence.",
    "Compliance is part of product design.",
    "Small controlled tests beat large uncontrolled bets.",
    "Do not confuse activity with progress.",
    "External rejection is information, not identity.",
    "Shipping creates opportunities that planning alone cannot.",
]
col1 = box(s, 0.9, 1.95, 5.8, 4.2)
col2 = box(s, 6.9, 1.95, 5.8, 4.2)
for i, t in enumerate(lessons):
    tgt = col1 if i < 5 else col2
    para(tgt, f"{i+1}.  {t}", 17, WHITE if i % 2 == 0 else LAV2, bold=False, first=(i in (0, 5)), space_after=14)
bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.7), Inches(6.3), Inches(10.0), Inches(0.85))
bar.fill.solid(); bar.fill.fore_color.rgb = PURPLE
bar.line.fill.background(); bar.shadow.inherit = False
tf = box(s, 2.0, 6.47, 9.4, 0.6)
para(tf, "\u201cYour background does not disqualify you. It gives you a different set of raw materials.\u201d", 17, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
footer(s, light=True)
notes(s, "Do not read all ten — pick the three that fit the room and tell them why.\n\nFor undergrads: #1 (no permission needed), #9 (rejection is information), #10 (shipping beats planning).\nFor MBA students: #5 (evidence over confidence), #7 (small controlled tests), #8 (activity vs. progress).\nFor MIS students: #4 (learn enough tech to ask better questions), #6 (compliance is product design).\n\nClose on the highlighted quote — it's the thesis of the entire talk.\n\n[20-MIN VERSION] Three lessons, 90 seconds.\n[45-60 MIN] Invite students to challenge one lesson they disagree with — productive debate starter.")

# ================= SLIDE 12 — WHAT COMES NEXT =================
s = slide(WHITE)
header(s, "What Comes Next", "The Next Chapter")
chip(s, 0.9, 1.9, 6.2, 0.55, GOLD, "CURRENT ROADMAP — NOT COMPLETED MILESTONES", 15, tcolor=DARK)
roadmap = [
    "Complete Google Play identity and organization verification",
    "Complete Android signing and first production build; test on real hardware",
    "Submit Dime Time to Google Play",
    "Receive Stripe Financial Connections approval → enable controlled public bank linking",
    "Receive Plaid Liabilities approval → verify automatic debt import",
    "Begin controlled public user acquisition",
    "Measure activation, retention, ACH volume and subscription conversion",
    "Raise $150,000 in pre-seed capital after launch milestones strengthen the story",
]
tf = box(s, 0.9, 2.7, 11.8, 4.0)
for i, t in enumerate(roadmap):
    para(tf, "→  " + t, 18, DARK, bold=(i >= 5), first=(i == 0), space_after=11)
tf = box(s, 0.9, 6.55, 11.8, 0.85)
para(tf, "No guaranteed approval dates or launch dates — external review timelines are not ours to promise.", 14, DEEP, bold=True, first=True, space_after=2)
para(tf, "Partner names indicate service relationships, not endorsements.", 14, DEEP)
footer(s)
notes(s, "Keep this factual and label it exactly as the banner says: a roadmap, not achievements.\n\nThe near-term sequence is mechanical: finish Google Play verification, Android build and submission, then the two external approvals — Stripe Financial Connections for public bank linking, Plaid Liabilities for automatic debt import. Those gates are outside our control, which is itself a lesson about regulated products.\n\nThe fundraise appears exactly once and last, on purpose: this is an academic talk, not a pitch. If asked, the plan is $150,000 pre-seed AFTER launch milestones strengthen the story.\n\n[DO NOT OVERSTATE] No dates, no 'expected approval by' promises.\n[20-MIN VERSION] 30 seconds — read the first and last items.")

# ================= SLIDE 13 — CLOSING =================
s = slide(DEEP)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.97), Inches(0.55), height=Inches(1.4))
tf = box(s, 0.8, 2.2, 11.7, 1.1)
para(tf, "WHAT WILL YOU BUILD?", 48, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
tf = box(s, 1.5, 3.6, 10.3, 1.3)
para(tf, "\u201cDime Time began as an idea about spare change. It became a lesson in persistence, systems thinking and turning education into execution.\u201d", 19, LAV2, align=PP_ALIGN.CENTER, first=True)
tf = box(s, 0.8, 5.15, 11.7, 1.2)
para(tf, "Tim Carlisle  —  Founder & CEO, Dime Time Technologies LLC", 18, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=4)
para(tf, "dime-time.com   •   tim@dime-time.com", 17, LAV2, align=PP_ALIGN.CENTER)
tf = box(s, 0.8, 6.55, 11.7, 0.6)
para(tf, "Get out of debt, one dime at a time with Dime Time!", 18, PURPLE, bold=True, align=PP_ALIGN.CENTER, first=True)
footer(s, light=True)
notes(s, "End with energy and turn it back to them: 'What will you build?' is a genuine question, not a rhetorical flourish.\n\nSuggested close: 'A year ago this was an idea in a folder. Today it's a live app with proven financial rails. The gap between those two states wasn't talent or money — it was starting, and then refusing to stop. Some of you have the idea already. What will you build?'\n\nThen open Q&A. Common questions to be ready for: How much did it cost? (development costs were kept minimal as a solo founder), How long did Apple review take? (multiple cycles), Are you hiring? (not yet — post-funding), Can I try the app? (yes — it's on the App Store today).\n\n[Reserve 15–20 minutes for Q&A in a 60-minute class session.]")

# ================= SLIDE 14 — APPENDIX — DISCUSSION QUESTIONS =================
s = slide(WHITE)
header(s, "Appendix — For Class Discussion", "Classroom Discussion Questions")
qs = [
    "Which part of Dime Time's moat is technical, and which is operational?",
    "When should a founder stop building and begin acquiring users?",
    "How should a fintech balance speed with safety?",
    "What should AI be trusted to do?",
    "What evidence would you require before investing?",
    "Which acquisition channel should Dime Time test first?",
    "How would you define an activated user?",
    "What ethical obligations come with automating financial behavior?",
]
tf = box(s, 0.9, 2.05, 11.8, 4.8)
for i, q in enumerate(qs):
    para(tf, f"{i+1}.  {q}", 19, DARK, first=(i == 0), space_after=14)
footer(s)
notes(s, "Optional slide for instructors — these work as small-group discussion prompts, exam questions, or a full class debate.\n\nStrongest for MBA rooms: #2 (build vs. acquire), #5 (evidence before investing) and #7 (defining activation).\nStrongest for MIS rooms: #1 (technical vs. operational moat) and #3 (speed vs. safety).\nStrongest for ethics-focused sessions: #8 — automating financial behavior carries real obligations, and students take it in surprising directions.\n\nOffer to stay after class or connect by email (tim@dime-time.com) for students who want to go deeper.")

assert _n == TOTAL, f"Slide count mismatch: built {_n}, expected {TOTAL}"
out = os.path.join(HERE, "dime-time-university-founder-journey.pptx")
prs.save(out)
print(f"Saved: {out} ({_n} slides)")
