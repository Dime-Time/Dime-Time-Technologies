"""Dime Time — 12-Month Timeline Deck (investor + academia)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

PURPLE = RGBColor(0x91, 0x8E, 0xF4)
DEEP = RGBColor(0x4A, 0x46, 0x8F)
DARK = RGBColor(0x2D, 0x2A, 0x4A)
LIGHT = RGBColor(0xF5, 0xF4, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xD6, 0x3A, 0x3A)
GREEN = RGBColor(0x2E, 0x9E, 0x62)
GOLD = RGBColor(0xE8, 0xB8, 0x4B)
LAV2 = RGBColor(0xC9, 0xC7, 0xF9)

LOGO = os.path.join(os.path.dirname(__file__), "..", "patent-deck-slides", "_logo_transparent.png")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def slide(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT, first=False, space_after=6, font="Poppins"):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font
    return p


def chip(s, x, y, w, h, fill, text, size, tcolor=WHITE, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    c = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = tcolor; r.font.name = "Poppins"
    return c


def month_header(s, kicker, title, commits, accent=PURPLE, tsize=40):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.28))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = box(s, 0.7, 0.55, 9.4, 0.5)
    para(tf, kicker.upper(), 15, DEEP, bold=True, first=True)
    tf2 = box(s, 0.7, 1.0, 9.6, 1.1)
    para(tf2, title, tsize, DARK, bold=True, first=True)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.7), Inches(0.6), Inches(2.0), Inches(2.0))
    c.fill.solid(); c.fill.fore_color.rgb = accent
    c.line.fill.background(); c.shadow.inherit = False
    ctf = c.text_frame; ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = ctf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(commits)
    r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Poppins"
    p2 = ctf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "commits"
    r2.font.size = Pt(13); r2.font.color.rgb = WHITE; r2.font.name = "Poppins"


def bullets(s, items, x=0.75, y=2.45, w=11.8, size=20, gap=12):
    tf = box(s, x, y, w, 4.6)
    for i, (t, emph) in enumerate(items):
        p = para(tf, t, size, DARK, bold=emph, first=(i == 0), space_after=gap)


# ---------- 1. TITLE ----------
s = slide(PURPLE)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.67), Inches(0.75), height=Inches(2.0))
tf = box(s, 1.0, 3.0, 11.3, 1.3)
para(tf, "DIME TIME", 60, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
tf = box(s, 1.0, 4.15, 11.3, 0.9)
para(tf, "One Founder. 12 Months. From Zero Code to Live Money Movement.", 26, DARK, bold=True, align=PP_ALIGN.CENTER, first=True)
tf = box(s, 1.0, 5.5, 11.3, 1.2)
para(tf, "Tim Carlisle, Founder", 20, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
para(tf, "U.S. Navy Veteran  |  MBA + MIS (UNLV)  |  Dallas, Texas", 16, WHITE, align=PP_ALIGN.CENTER)

# ---------- 2. IDEA ERA ----------
s = slide(LIGHT)
tf = box(s, 0.7, 0.5, 12, 0.6)
para(tf, "THE IDEA ERA  •  2014–2025", 15, DEEP, bold=True, first=True)
tf = box(s, 1.0, 1.15, 11.3, 1.5)
para(tf, "11 YEARS", 88, PURPLE, bold=True, align=PP_ALIGN.CENTER, first=True)
tf = box(s, 1.0, 2.85, 11.3, 1.9)
para(tf, "No Funding", 34, DARK, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=4)
para(tf, "No Team", 34, DARK, bold=True, align=PP_ALIGN.CENTER, space_after=4)
para(tf, "No Code", 34, DARK, bold=True, align=PP_ALIGN.CENTER)
tf = box(s, 1.0, 5.1, 11.3, 0.6)
para(tf, "Just an idea: redirect spare change toward DEBT, not savings.", 19, DEEP, align=PP_ALIGN.CENTER, first=True)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.42), Inches(5.85), Inches(2.5), Pt(2.5))
ln.fill.solid(); ln.fill.fore_color.rgb = PURPLE; ln.line.fill.background(); ln.shadow.inherit = False
tf = box(s, 1.0, 6.05, 11.3, 1.1)
para(tf, "August 3, 2025", 28, DARK, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=2)
para(tf, "Everything changed.", 22, PURPLE, bold=True, align=PP_ALIGN.CENTER)

# ---------- 3. YEAR AT A GLANCE ----------
s = slide(WHITE)
tf = box(s, 0.7, 0.4, 12, 0.6)
para(tf, "THE YEAR AT A GLANCE", 15, DEEP, bold=True, first=True)
tf = box(s, 0.7, 0.85, 12, 0.8)
para(tf, "1,099 commits, month by month", 34, DARK, bold=True, first=True)
months = [("Aug", 106), ("Sep", 206), ("Oct", 105), ("Nov", 127), ("Dec", 140), ("Jan", 21),
          ("Feb", 13), ("Mar", 31), ("Apr", 39), ("May", 69), ("Jun", 105), ("Jul", 137)]
chart_x, chart_y, chart_w, chart_h = 0.9, 2.1, 11.6, 3.4
bw = chart_w / len(months) * 0.62
gapx = chart_w / len(months)
maxv = 206
for i, (m, v) in enumerate(months):
    bh = chart_h * v / maxv
    x = chart_x + i * gapx + (gapx - bw) / 2
    y = chart_y + chart_h - bh
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(bw), Inches(bh))
    b.fill.solid(); b.fill.fore_color.rgb = PURPLE if v >= 100 else LAV2
    b.line.fill.background(); b.shadow.inherit = False
    tf = box(s, x - 0.2, y - 0.42, bw + 0.4, 0.4)
    para(tf, str(v), 13, DARK, bold=True, align=PP_ALIGN.CENTER, first=True)
    tf = box(s, x - 0.25, chart_y + chart_h + 0.06, bw + 0.5, 0.35)
    para(tf, m, 13, DEEP, bold=True, align=PP_ALIGN.CENTER, first=True)
marks = [(0, "First Commit", DEEP), (1, "Rejected 9/19", RED), (7, "Mercury", GOLD),
         (9, "Stripe", GREEN), (10, "App Store LIVE", PURPLE), (11, "First Real $1", DARK)]
for mi, label, mc in marks:
    x = chart_x + mi * gapx + gapx / 2
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.09), Inches(chart_y + chart_h + 0.44), Inches(0.18), Inches(0.18))
    d.fill.solid(); d.fill.fore_color.rgb = mc
    d.line.fill.background(); d.shadow.inherit = False
lg_x = 0.9
for mi, label, mc in marks:
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(lg_x), Inches(6.62), Inches(0.16), Inches(0.16))
    d.fill.solid(); d.fill.fore_color.rgb = mc
    d.line.fill.background(); d.shadow.inherit = False
    tf = box(s, lg_x + 0.22, 6.5, 1.85, 0.4)
    para(tf, label, 12, DARK, bold=True, first=True)
    lg_x += 0.35 + len(label) * 0.088 + 0.35

# ---------- 4. MONTH 1 ----------
s = slide(LIGHT)
month_header(s, "Month 1  •  August 2025", "The Explosion", 106)
bullets(s, [
    ("First commit: August 3, 2025.", True),
    ("Core app, purple brand, and crypto round-ups — built in the FIRST WEEKEND.", False),
    ("By August 12: database, bank linking, PIN security, and the round-up engine.", False),
    ("Day 11: iPhone app development begins.", False),
    ("A first iOS app (\u201cDime Time Technologies\u201d) starts taking shape — ~60 commits of work that won't survive the year.", False),
])
tf = box(s, 0.75, 6.5, 11.8, 0.6)
para(tf, "Tip for the live talk: drop an early app screenshot here — primitive beginnings make the finished product hit harder.", 13, DEEP, first=True)

# ---------- 5. MONTH 2 / SEPT 19 ----------
s = slide(WHITE)
month_header(s, "Month 2  •  September 2025", "Biggest Month. First Heartbreak.", 206, accent=RED, tsize=34)
tf = box(s, 0.75, 2.3, 6.6, 3.6)
para(tf, "Lavender brand locked in. Legal agreements. Analytics.", 18, DARK, first=True, space_after=10)
para(tf, "Sila Money wired up as the ACH test sandbox — then quoted $5,000/month. Never viable.", 18, DARK, space_after=10)
para(tf, "TestFlight, CI/CD pipeline, and encryption all arrive the same month.", 18, DARK)
pnl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(2.2), Inches(5.0), Inches(4.3))
pnl.fill.solid(); pnl.fill.fore_color.rgb = LIGHT
pnl.line.color.rgb = RED; pnl.line.width = Pt(2); pnl.shadow.inherit = False
tf = box(s, 7.9, 2.5, 4.4, 3.9)
para(tf, "SEPTEMBER 19", 30, DARK, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=2)
para(tf, "Founder's Birthday", 18, DEEP, align=PP_ALIGN.CENTER, space_after=2)
para(tf, "First App Submitted to Apple", 18, DEEP, align=PP_ALIGN.CENTER, space_after=10)
para(tf, "REJECTED", 54, RED, bold=True, align=PP_ALIGN.CENTER, space_after=8)
para(tf, "TestFlight opened to a blank landing page. Weeks of fixes couldn't save it.", 14, DARK, align=PP_ALIGN.CENTER)

# ---------- 6. MONTHS 3–5 ----------
s = slide(DEEP)
tf = box(s, 0.7, 0.6, 12, 0.6)
para(tf, "MONTHS 3–5  •  OCTOBER – DECEMBER 2025", 15, LAV2, bold=True, first=True)
tf = box(s, 1.0, 1.5, 11.3, 1.6)
para(tf, "372 COMMITS", 76, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
tf = box(s, 1.0, 3.2, 11.3, 1.6)
para(tf, "Three Months. Complete Rebuild.", 32, PURPLE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=14)
para(tf, "First app scrapped. iOS project reborn under a new identity.", 20, WHITE, align=PP_ALIGN.CENTER)
tf = box(s, 1.0, 5.1, 11.3, 1.7)
para(tf, "Certificates. Provisioning profiles. Signed builds. The exact machinery that broke the first app — rebuilt properly this time.", 18, LAV2, align=PP_ALIGN.CENTER, first=True, space_after=10)
para(tf, "December's log: \u201cBuild application archive and export successfully.\u201d", 18, GOLD, bold=True, align=PP_ALIGN.CENTER)

# ---------- 7. MONTHS 6–7 ----------
s = slide(LIGHT)
month_header(s, "Months 6–7  •  January – February 2026", "The Quiet Grind", 34)
bullets(s, [
    ("Lowest commit counts of the year. Highest-stakes work of the year.", True),
    ("Face ID + PIN app lock. Encrypted tokens. Hardened passwords. Rate limiting.", False),
    ("The unglamorous security work fintechs live or die on.", False),
    ("February 21: the first patent documentation enters the project.", False),
    ("All of it built nights and weekends around a 30–40 hour/week sommelier job.", False),
])

# ---------- 8. MONTH 8 ----------
s = slide(WHITE)
month_header(s, "Month 8  •  March 2026", "The Bank That Said Yes", 31)
rows = [("✗  Axos", RED), ("✗  Brex", RED), ("✗  Two more banks", RED), ("✓  MERCURY", GREEN)]
y = 2.35
for t, c in rows:
    chip(s, 0.9, y, 5.4, 0.85, WHITE if c == RED else GREEN, t, 24, tcolor=c if c == RED else WHITE)
    if c == RED:
        sh = s.shapes[-1]; sh.line.color.rgb = RED; sh.line.width = Pt(1.5)
    y += 1.05
tf = box(s, 6.9, 2.5, 5.7, 3.6)
para(tf, "Business banking was a wall of rejections — until Mercury said yes.", 22, DARK, bold=True, first=True, space_after=12)
para(tf, "March 28: Mercury becomes Dime Time's business bank. Sila's sandbox era ends. Plaid moves to production.", 18, DARK, space_after=12)
para(tf, "The money infrastructure takes its final shape.", 18, DEEP, bold=True)

# ---------- 9. MONTHS 9–10 ----------
s = slide(LIGHT)
month_header(s, "Months 9–10  •  April – May 2026", "Finding the Rails", 108)
p1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.3), Inches(5.5), Inches(2.1))
p1.fill.solid(); p1.fill.fore_color.rgb = WHITE
p1.line.color.rgb = RED; p1.line.width = Pt(2); p1.shadow.inherit = False
tf = box(s, 1.2, 2.55, 4.9, 1.7)
para(tf, "Plaid said no.", 30, RED, bold=True, first=True, space_after=6)
para(tf, "Rejected for ACH transfers. Dwolla and others evaluated next.", 16, DARK)
p2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(2.3), Inches(5.5), Inches(2.1))
p2.fill.solid(); p2.fill.fore_color.rgb = GREEN
p2.line.fill.background(); p2.shadow.inherit = False
tf = box(s, 7.2, 2.55, 4.9, 1.7)
para(tf, "Stripe said yes.", 30, WHITE, bold=True, first=True, space_after=6)
para(tf, "May 27, 2026: payment rails secured. Money movement had a home.", 16, WHITE)
tf = box(s, 0.9, 4.75, 11.6, 2.0)
para(tf, "Also this stretch:  D-U-N-S number  •  provisional patent drafted (USPTO-ready)  •  patent deck + business plan  •  new marketing site with legal pages", 17, DARK, first=True, space_after=8)
para(tf, "And a security scare handled by the book — an internal config slip exposed an encryption key inside the project (never externally). Treated as compromised anyway: new key, every credential re-encrypted, zero user impact.", 15, DEEP)

# ---------- 10. JUNE 29 ----------
s = slide(PURPLE)
tf = box(s, 0.7, 0.55, 12, 0.6)
para(tf, "MONTH 11  •  JUNE 2026  •  105 COMMITS", 15, WHITE, bold=True, first=True)
tf = box(s, 1.0, 1.6, 11.3, 1.8)
para(tf, "JUNE 29", 96, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
tf = box(s, 1.0, 3.6, 11.3, 0.9)
para(tf, "DIME TIME GOES LIVE ON THE APPLE APP STORE", 28, DARK, bold=True, align=PP_ALIGN.CENTER, first=True)
tf = box(s, 1.0, 4.8, 11.3, 1.9)
para(tf, "Apple's last punch — a \u201cduplicate app\u201d rejection caused by the first app's ghost record — cleared through the Resolution Center.", 17, WHITE, align=PP_ALIGN.CENTER, first=True, space_after=10)
para(tf, "283 days after the birthday rejection, the answer flipped to yes.", 20, DARK, bold=True, align=PP_ALIGN.CENTER)

# ---------- 11. REAL MONEY FLOW ----------
s = slide(WHITE)
month_header(s, "Month 12  •  July 2026", "Real Money", 137)
steps = [("$1.00", "Real ACH debit — the founder's own account"),
         ("CHASE", "Pulled from a real consumer bank"),
         ("STRIPE", "Processed on live payment rails"),
         ("MERCURY", "Settled into the business bank"),
         ("PROVEN", "\u201cDIME TIME LLC\u201d on a real bank statement")]
x = 0.7
bw2 = 2.25; gap2 = 0.32
for i, (big, small) in enumerate(steps):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5), Inches(bw2), Inches(1.9))
    c.fill.solid(); c.fill.fore_color.rgb = PURPLE if big != "PROVEN" else GREEN
    c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = big
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Poppins"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = small
    r2.font.size = Pt(10.5); r2.font.color.rgb = WHITE; r2.font.name = "Poppins"
    if i < len(steps) - 1:
        tf2 = box(s, x + bw2 - 0.06, 3.15, gap2 + 0.14, 0.6)
        para(tf2, "→", 24, DARK, bold=True, align=PP_ALIGN.CENTER, first=True)
    x += bw2 + gap2
tf = box(s, 0.75, 5.0, 11.8, 1.9)
para(tf, "July 7–8: initiated  •  July 13–14: settled  •  July 15: payout landed in Mercury", 18, DEEP, bold=True, first=True, space_after=10)
para(tf, "Also this month: v1.0.4 redesign approved by Apple, admin money controls, automatic debt import built, AI-readable web guides.", 16, DARK, space_after=10)
para(tf, "The machine works — proven with real money, 346 days after the first commit.", 18, DARK, bold=True)

# ---------- 12. WHY DIME TIME EXISTS ----------
s = slide(DEEP)
tf = box(s, 0.7, 0.6, 12, 0.6)
para(tf, "WHY DIME TIME EXISTS", 15, LAV2, bold=True, first=True)
tf = box(s, 1.2, 1.7, 10.9, 4.5)
para(tf, "Americans carry trillions of dollars of consumer debt.", 34, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=22)
para(tf, "Most financial apps help people save.", 26, LAV2, align=PP_ALIGN.CENTER, space_after=22)
para(tf, "Dime Time was built to help people owe less.", 34, PURPLE, bold=True, align=PP_ALIGN.CENTER)

# ---------- 13. THE GAUNTLET ----------
s = slide(LIGHT)
tf = box(s, 0.7, 0.5, 12, 0.6)
para(tf, "THE GAUNTLET", 15, DEEP, bold=True, first=True)
tf = box(s, 0.7, 0.95, 12, 0.8)
para(tf, "Every rejection was a redirect.", 34, DARK, bold=True, first=True)
c1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.1), Inches(5.5), Inches(4.6))
c1.fill.solid(); c1.fill.fore_color.rgb = WHITE
c1.line.color.rgb = RED; c1.line.width = Pt(2); c1.shadow.inherit = False
tf = box(s, 1.25, 2.35, 4.8, 4.1)
para(tf, "NO", 40, RED, bold=True, first=True, space_after=10)
for t in ["Apple — birthday submission", "Sila — $5,000/month pricing", "Plaid — ACH transfers", "Axos", "Brex", "Two more banks"]:
    para(tf, "✗  " + t, 19, DARK, space_after=8)
c2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(2.1), Inches(5.5), Inches(4.6))
c2.fill.solid(); c2.fill.fore_color.rgb = GREEN
c2.line.fill.background(); c2.shadow.inherit = False
tf = box(s, 7.25, 2.35, 4.8, 4.1)
para(tf, "YES", 40, WHITE, bold=True, first=True, space_after=10)
for t in ["Mercury — business banking", "Stripe — payment rails", "Apple — June 29, 2026", "USPTO — provisional patent drafted", "The first real dollar — July 2026"]:
    para(tf, "✓  " + t, 19, WHITE, space_after=8)

# ---------- 14. STOREFRONT ----------
s = slide(WHITE)
tf = box(s, 0.7, 0.5, 12, 0.6)
para(tf, "DIME-TIME.COM", 15, DEEP, bold=True, first=True)
tf = box(s, 0.7, 0.95, 12, 0.8)
para(tf, "The Storefront Before the Store", 34, DARK, bold=True, first=True)
items = [
    ("Aug 2025", "Website live with the lavender brand — two weeks after the first commit"),
    ("Sep 2025", "Privacy policy, Google Analytics, veteran-owned badges"),
    ("Oct 2025", "QR codes handed out at a finance conference — before the app existed"),
    ("May 2026", "Full marketing site: FAQ, legal pages, bot-protected contact form"),
    ("Jul 2026", "SEO + AI-readable guides — built to be found by Google AND ChatGPT"),
]
y = 2.15
for d, t in items:
    chip(s, 0.9, y, 1.7, 0.68, PURPLE, d, 15)
    tf = box(s, 2.85, y + 0.05, 9.6, 0.7)
    para(tf, t, 18, DARK, first=True)
    y += 0.88
tf = box(s, 0.9, 6.7, 11.6, 0.6)
para(tf, "The storefront existed before the store — and it was already measuring visitors.", 16, DEEP, bold=True, first=True)

# ---------- 15. BY THE NUMBERS ----------
s = slide(PURPLE)
tf = box(s, 0.7, 0.5, 12, 0.6)
para(tf, "BY THE NUMBERS", 15, WHITE, bold=True, first=True)
stats = [("1,160+", "commits"), ("35,000+", "lines of code"), ("12", "months"),
         ("1", "founder"), ("$0", "raised"), ("1", "live fintech")]
positions = [(0.9, 1.6), (5.05, 1.6), (9.2, 1.6), (0.9, 4.35), (5.05, 4.35), (9.2, 4.35)]
for (bigt, small), (x, y) in zip(stats, positions):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.3), Inches(2.3))
    c.fill.solid(); c.fill.fore_color.rgb = WHITE
    c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = bigt
    r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name = "Poppins"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = small
    r2.font.size = Pt(18); r2.font.color.rgb = DARK; r2.font.name = "Poppins"
tf = box(s, 0.9, 6.85, 11.6, 0.5)
para(tf, "Built while working 30–40 hours a week on his feet.", 16, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)

# ---------- 16. ROADMAP ----------
s = slide(LIGHT)
tf = box(s, 0.7, 0.5, 12, 0.6)
para(tf, "WHAT'S NEXT", 15, DEEP, bold=True, first=True)
tf = box(s, 0.7, 0.95, 12, 0.8)
para(tf, "The road from here is measured in users.", 34, DARK, bold=True, first=True)
steps2 = ["100 Users", "1,000 Users", "National Launch", "Institution Partnerships", "Scale"]
x = 0.7; bw3 = 2.3; gap3 = 0.28
for i, t in enumerate(steps2):
    shade = [LAV2, LAV2, PURPLE, PURPLE, DEEP][i]
    chip(s, x, 3.0, bw3, 1.4, shade, t, 18, tcolor=(DARK if shade == LAV2 else WHITE))
    if i < len(steps2) - 1:
        tf = box(s, x + bw3 - 0.06, 3.35, gap3 + 0.14, 0.6)
        para(tf, "→", 22, DARK, bold=True, align=PP_ALIGN.CENTER, first=True)
    x += bw3 + gap3
tf = box(s, 0.9, 5.1, 11.6, 1.6)
para(tf, "Near-term: automatic debt import in production  •  first paying subscribers ($2.99/mo)  •  the public launch announcement.", 18, DARK, first=True, space_after=8)
para(tf, "Positioning: a case study in building a regulated fintech from scratch — while working full-time.", 17, DEEP, bold=True)

# ---------- 17. CLOSING ----------
s = slide(DEEP)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.92), Inches(0.7), height=Inches(1.5))
tf = box(s, 1.0, 2.5, 11.3, 3.4)
para(tf, "Built by one person.", 34, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=6)
para(tf, "Nights.", 34, LAV2, bold=True, align=PP_ALIGN.CENTER, space_after=6)
para(tf, "Weekends.", 34, LAV2, bold=True, align=PP_ALIGN.CENTER, space_after=6)
para(tf, "Seven rejections.", 34, WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=6)
para(tf, "Still launched.", 40, PURPLE, bold=True, align=PP_ALIGN.CENTER)
tf = box(s, 1.0, 6.3, 11.3, 0.7)
para(tf, "dime-time.com   |   Available on the App Store", 18, WHITE, align=PP_ALIGN.CENTER, first=True)

out = os.path.join(os.path.dirname(__file__), "dime-time-timeline-deck.pptx")
prs.save(out)
print("Saved:", out)
