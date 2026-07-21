"""Dime Time investor pitch deck — full 14-slide generator (July 2026 revision).

Builds the entire deck natively in python-pptx (searchable/selectable text in
the exported PDF, sequential page numbers). Replaces the old PNG-stack build.

Slides: 01 cover, 02 problem, 03 solution, 04 allocation, 05 product,
06 market, 07 business model, 08 competition, 09 defensibility,
10 go-to-market, 11 traction, 12 team, 13 raise, 14 company info.

Variants: dime-time-pitch-deck.pptx (14) and -with-ip.pptx (14 + 3 IP PNGs).

Verified stats baked in (do not change without re-verifying):
- $18.8T U.S. household debt / $1.25T credit card — NY Fed Household Debt &
  Credit Report Q1 2026, released May 12, 2026.
- 91% of U.S. adults own a smartphone — Pew Research Center, Mobile Fact
  Sheet, 2025 survey.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

HERE = os.path.dirname(os.path.abspath(__file__))

LAVENDER = RGBColor(0x91, 0x8E, 0xF4)
LAV_PANEL = RGBColor(0xA4, 0xA1, 0xF6)
DARKBG = RGBColor(0x21, 0x1D, 0x3D)
CARD = RGBColor(0x2B, 0x27, 0x4E)
CARD_EDGE = RGBColor(0x4A, 0x46, 0x8F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FAINT = RGBColor(0xC9, 0xC7, 0xF9)
DEEP = RGBColor(0x4A, 0x46, 0x8F)
MUTED_DARK = RGBColor(0x9B, 0x98, 0xC4)   # small text on dark bg
LAV_SOFT = RGBColor(0xE9, 0xE8, 0xFD)     # small text on lavender bg
GREEN = RGBColor(0x34, 0xC7, 0x8A)

TOTAL = 14

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def slide(bg):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    for i, (t, size, color, bold, sp) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tf


def spaced(t):
    return " ".join(list(t.replace(" ", "  ")))


def footer(s, n, light):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(7.06), Inches(0.42), Pt(1.6))
    ln.fill.solid(); ln.fill.fore_color.rgb = WHITE
    ln.line.fill.background(); ln.shadow.inherit = False
    text(s, 1.08, 6.9, 3.0, 0.35, [(spaced("DIME TIME"), 11, WHITE, True, 0)])
    num_c = LAV_SOFT if light else FAINT
    text(s, 11.6, 6.95, 1.3, 0.35, [(f"{n:02d} / {TOTAL}", 11, num_c, False, 0)], align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, dark=True, edge=None):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.adjustments[0] = min(0.12, 0.35 / min(w, h))
    if dark:
        c.fill.solid(); c.fill.fore_color.rgb = CARD
        c.line.color.rgb = edge or CARD_EDGE; c.line.width = Pt(1)
    else:
        c.fill.solid(); c.fill.fore_color.rgb = LAV_PANEL
        c.line.fill.background()
    c.shadow.inherit = False
    return c


def badge(s, x, y, label, w=1.55, filled=True):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    b.adjustments[0] = 0.5
    if filled:
        b.fill.solid(); b.fill.fore_color.rgb = DEEP
        b.line.fill.background()
    else:
        b.fill.background()
        b.line.color.rgb = FAINT; b.line.width = Pt(1)
    b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = spaced(label)
    r.font.size = Pt(8.5); r.font.bold = True; r.font.color.rgb = WHITE
    return b


# ---------- iPhone framing ----------
SHOTS = os.path.join(HERE, "..", "Dime-Time-App-Store-Screenshots-FINAL", "iPhone-6.9")
FRAMED_DIR = "/tmp/framed-phones"


def build_framed(src, dst):
    from PIL import Image, ImageDraw
    shot = Image.open(src).convert("RGB")
    Wp, Hp = shot.size
    bez, rad_screen = 40, 150
    rad_body = rad_screen + bez
    body_w, body_h = Wp + 2 * bez, Hp + 2 * bez
    btn = 16
    canvas = Image.new("RGBA", (body_w + 2 * btn, body_h), (0, 0, 0, 0))
    d = ImageDraw.Draw(canvas)
    BODY, BTN = (22, 20, 42, 255), (38, 36, 62, 255)
    for y0, y1 in [(560, 700), (790, 1010), (1080, 1300)]:
        d.rounded_rectangle([4, y0, btn + 8, y1], 8, fill=BTN)
    d.rounded_rectangle([canvas.width - btn - 8, 880, canvas.width - 4, 1230], 8, fill=BTN)
    d.rounded_rectangle([btn, 0, btn + body_w - 1, body_h - 1], rad_body, fill=BODY)
    mask = Image.new("L", (Wp, Hp), 0)
    ImageDraw.Draw(mask).rounded_rectangle([0, 0, Wp - 1, Hp - 1], rad_screen, fill=255)
    canvas.paste(shot, (btn + bez, bez), mask)
    iw, ih = 380, 110
    ix, iy = btn + bez + (Wp - iw) // 2, bez + 52
    d.rounded_rectangle([ix, iy, ix + iw, iy + ih], ih // 2, fill=(8, 7, 18, 255))
    canvas.save(dst)
    return canvas.width / canvas.height


def framed(name):
    os.makedirs(FRAMED_DIR, exist_ok=True)
    dst = os.path.join(FRAMED_DIR, name)
    aspect = build_framed(os.path.join(SHOTS, name), dst)
    return dst, aspect


def phone(s, x, y, h, name):
    dst, aspect = framed(name)
    w = h * aspect
    s.shapes.add_picture(dst, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    return w


def glyph_png():
    """White alarm-clock glyph from the app icon, transparent background."""
    from PIL import Image
    src = os.path.join(HERE, "..", "..", "client", "src", "assets", "dime-time-app-icon.png")
    dst = "/tmp/dt-glyph.png"
    im = Image.open(src).convert("RGBA")
    px = im.load()
    for yy in range(im.height):
        for xx in range(im.width):
            r, g, b, a = px[xx, yy]
            if r > 215 and g > 215 and b > 235:
                px[xx, yy] = (255, 255, 255, a)
            else:
                px[xx, yy] = (255, 255, 255, 0)
    im.save(dst)
    return dst


# ================= SLIDE 01 — COVER =================
s = slide(LAVENDER)
gw = 1.5
s.shapes.add_picture(glyph_png(), Inches((13.333 - gw) / 2), Inches(1.0), width=Inches(gw), height=Inches(gw))
text(s, 1.0, 2.55, 11.333, 1.1, [(spaced("DIME TIME"), 54, WHITE, True, 0)], align=PP_ALIGN.CENTER)
text(s, 1.0, 3.8, 11.333, 0.5, [("Get out of debt, one dime at a time.", 19, WHITE, False, 0)], align=PP_ALIGN.CENTER)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.42), Inches(5.55), Inches(2.5), Pt(1.2))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xBD, 0xBB, 0xF8); ln.line.fill.background(); ln.shadow.inherit = False
text(s, 1.0, 5.75, 11.333, 0.4, [("Tim Carlisle · Founder & CEO", 14, WHITE, True, 0)], align=PP_ALIGN.CENTER)
text(s, 1.0, 6.2, 11.333, 0.4, [(spaced("AUTOMATED DEBT REDUCTION THROUGH EVERYDAY SPENDING"), 10.5, LAV_SOFT, False, 0)], align=PP_ALIGN.CENTER)

# ================= SLIDE 02 — PROBLEM =================
s = slide(DARKBG)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("THE PROBLEM"), 12, LAVENDER, True, 0)])
text(s, 0.62, 1.0, 11.5, 0.9, [("Americans Are Drowning in Debt", 34, WHITE, True, 0)])
text(s, 0.72, 2.1, 6.6, 1.3, [("$18.8T", 60, WHITE, True, 2),
                              (spaced("U.S. HOUSEHOLD DEBT"), 12, LAVENDER, True, 0)])
prob_bullets = [
    "$1.25T of it is credit-card debt — the segment Dime Time attacks first",
    "People intend to pay down debt — but don't follow through",
    "Manual extra payments require discipline most users don't maintain",
    "Financial stress is rising across all generations",
]
text(s, 0.78, 3.9, 6.9, 2.6, [("\u2022  " + b, 14.5, WHITE, False, 12) for b in prob_bullets])
card(s, 8.1, 2.1, 4.5, 4.0)
text(s, 8.45, 2.45, 3.9, 0.4, [(spaced("WHY NOW"), 11, LAVENDER, True, 0)])
text(s, 8.45, 2.95, 3.9, 1.2, [("91%", 44, WHITE, True, 2),
                               ("of U.S. adults own a smartphone", 13, FAINT, False, 0)])
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.45), Inches(4.75), Inches(3.8), Pt(1))
ln.fill.solid(); ln.fill.fore_color.rgb = CARD_EDGE; ln.line.fill.background(); ln.shadow.inherit = False
text(s, 8.45, 4.95, 3.9, 1.0, [("Plaid, ACH rails and mobile banking now make fully automated debt repayment possible.", 12.5, WHITE, False, 0)])
text(s, 0.72, 6.5, 11.9, 0.35, [("Sources: Federal Reserve Bank of New York, Household Debt & Credit Report Q1 2026 (released May 2026) · Pew Research Center, Mobile Fact Sheet (2025)", 9, MUTED_DARK, False, 0)])
footer(s, 2, light=False)

# ================= SLIDE 03 — SOLUTION =================
s = slide(LAVENDER)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("THE SOLUTION"), 12, WHITE, True, 0)])
text(s, 0.62, 1.0, 9.0, 0.9, [("Automated Debt Reduction", 34, WHITE, True, 0)])
sol_bullets = [
    "Every purchase is rounded up automatically",
    "Users control how their spare change is allocated",
    "Round-ups accumulate, then fund user-authorized ACH debt payments",
    "No willpower required — it works quietly in the background",
]
text(s, 0.95, 2.3, 8.0, 3.2, [("\u2022  " + b, 17, WHITE, False, 20) for b in sol_bullets])
text(s, 0.95, 5.6, 8.0, 0.8, [("Round-ups are batched and transferred on a schedule the user approves — money does not move on every individual purchase.", 12, LAV_SOFT, False, 0)])
phone(s, 9.75, 1.35, 5.35, "01-dashboard-iphone.png")
footer(s, 3, light=True)

# ================= SLIDE 04 — DYNAMIC ALLOCATION =================
s = slide(LAVENDER)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("THE BREAKTHROUGH"), 12, WHITE, True, 0)])
text(s, 0.62, 1.0, 11.5, 0.9, [("Dynamic Allocation Engine", 34, WHITE, True, 0)])
text(s, 0.66, 1.78, 9.0, 0.4, [("Automated allocation from every transaction.", 15, WHITE, True, 0)])
pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(2.75), Inches(2.95), Inches(1.15))
pill.adjustments[0] = 0.5
pill.fill.solid(); pill.fill.fore_color.rgb = WHITE; pill.line.fill.background(); pill.shadow.inherit = False
text(s, 0.98, 2.95, 2.65, 0.3, [(spaced("ROUND-UP EXAMPLE"), 8.5, DEEP, True, 0)])
text(s, 0.98, 3.22, 2.5, 0.5, [("$0.37", 26, DEEP, True, 0)])
for cy, name, sub, pct, blabel, bfill in [
    (2.35, "Debt Repayment", "primary debt payoff", "80%", "LIVE IN APP", True),
    (4.05, "Bitcoin (Coinbase)", "long-term wealth — not yet available to users", "20%", "PLANNED", False),
]:
    c = card(s, 7.0, cy, 5.6, 1.35, dark=False)
    text(s, 7.3, cy + 0.18, 3.6, 0.5, [(name, 17, WHITE, True, 2), (sub, 10.5, LAV_SOFT, False, 0)])
    text(s, 11.2, cy + 0.28, 1.2, 0.6, [(pct, 26, WHITE, True, 0)], align=PP_ALIGN.RIGHT)
    badge(s, 7.3, cy + 0.92, blabel, w=1.7, filled=bfill)
for x0, b in [(0.72, "Real-time allocation tracking"), (4.9, "User-controlled split"), (9.0, "Payments execute over ACH rails")]:
    text(s, x0, 5.85, 4.0, 0.6, [("\u2022  " + b, 13, WHITE, False, 0)])
text(s, 2.0, 6.55, 9.3, 0.4, [("Patent-pending allocation workflow · Real transfers roll out with public bank linking.", 12, LAV_SOFT, False, 0)], align=PP_ALIGN.CENTER)
footer(s, 4, light=True)

# ================= SLIDE 05 — PRODUCT =================
s = slide(DARKBG)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("THE PRODUCT"), 12, LAVENDER, True, 0)])
text(s, 0.62, 1.0, 11.5, 0.9, [("Live Product — Real Financial Movement", 32, WHITE, True, 0)])
shots = [
    ("01-dashboard-iphone.png", "Live dashboard"),
    ("02-debts-iphone.png", "Debt tracking"),
    ("03-roundups-iphone.png", "Round-ups"),
    ("04-insights-iphone.png", "Insights"),
]
ph_h = 3.85
ph_w = ph_h * framed(shots[0][0])[1]
gap = (13.333 - 2 * 1.15 - 4 * ph_w) / 3
x = 1.15
for png, label in shots:
    phone(s, x, 2.0, ph_h, png)
    text(s, x - 0.45, 6.0, ph_w + 0.9, 0.4, [(label, 14, WHITE, True, 0)], align=PP_ALIGN.CENTER)
    x += ph_w + gap
text(s, 2.0, 6.5, 9.3, 0.35, [("Actual v1.0.5 App Store screenshots — live on the Apple App Store, July 2026", 12, FAINT, False, 0)], align=PP_ALIGN.CENTER)
footer(s, 5, light=False)

# ================= SLIDE 06 — MARKET =================
s = slide(LAVENDER)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("THE MARKET"), 12, WHITE, True, 0)])
text(s, 0.62, 1.0, 11.8, 0.9, [("Market Opportunity, Built Bottom-Up", 34, WHITE, True, 0)])
# left: context
c = card(s, 0.72, 2.15, 4.7, 3.9, dark=False)
text(s, 1.05, 2.45, 4.0, 0.35, [(spaced("MARKET CONTEXT"), 11, WHITE, True, 0)])
text(s, 1.05, 2.9, 4.0, 1.9, [("$18.8T", 40, WHITE, True, 2),
                              ("U.S. household debt", 12.5, WHITE, False, 10),
                              ("$1.25T", 40, WHITE, True, 2),
                              ("credit-card balances — our first target", 12.5, WHITE, False, 0)])
text(s, 1.05, 5.5, 4.05, 0.5, [("Context for the size of the problem — not our revenue model.", 10.5, LAV_SOFT, False, 0)])
# right: bottom-up ARR
text(s, 5.95, 2.2, 6.7, 0.35, [(spaced("SUBSCRIPTION REVENUE AT $2.99 / MONTH"), 11, WHITE, True, 0)])
rows = [("1,400", "paying subscribers", "$50,232 ARR"),
        ("5,600", "paying subscribers", "$200,928 ARR"),
        ("25,000", "paying subscribers", "$897,000 ARR")]
y = 2.75
for n, lbl, arr in rows:
    text(s, 5.95, y, 1.8, 0.5, [(n, 26, WHITE, True, 0)])
    text(s, 7.85, y + 0.14, 2.2, 0.4, [(lbl, 12.5, LAV_SOFT, False, 0)])
    text(s, 9.65, y + 0.06, 3.0, 0.5, [(arr, 18, WHITE, True, 0)], align=PP_ALIGN.RIGHT)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.95), Inches(y + 0.62), Inches(6.7), Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xB8, 0xB6, 0xF8); ln.line.fill.background(); ln.shadow.inherit = False
    y += 0.92
text(s, 5.95, y + 0.05, 6.7, 0.7, [("Subscriber counts are illustrative milestones — targets, not forecasts. ARR = subscribers × $2.99 × 12.", 10.5, LAV_SOFT, False, 0)])
text(s, 0.72, 6.5, 11.9, 0.35, [("Source: Federal Reserve Bank of New York, Household Debt & Credit Report Q1 2026 (released May 2026)", 9, LAV_SOFT, False, 0)])
footer(s, 6, light=True)

# ================= SLIDE 07 — BUSINESS MODEL =================
s = slide(DARKBG)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("BUSINESS MODEL"), 12, LAVENDER, True, 0)])
text(s, 0.62, 1.0, 11.5, 0.9, [("Simple, Scalable Revenue Model", 34, WHITE, True, 0)])
tiers = [
    ("CORE", "$2.99", "BUILT · FEATURE-GATED", True,
     ["Round-up automation", "Debt-first allocation", "Progress tracking & insights"]),
    ("PRO", "$3.99", "PLANNED", False,
     ["Enhanced allocation controls", "Advanced analytics", "Deeper automation"]),
    ("PREMIUM", "$4.99", "PLANNED", False,
     ["Faster transfer options", "Multiple funding accounts", "Priority routing"]),
]
x = 0.72
for name, price, blabel, bfill, feats in tiers:
    card(s, x, 2.0, 3.86, 3.7)
    text(s, x + 0.32, 2.3, 3.2, 0.35, [(spaced(name), 13, WHITE, True, 0)])
    badge(s, x + 0.32, 2.72, blabel, w=3.1 if len(blabel) > 8 else 1.4, filled=bfill)
    text(s, x + 0.32, 3.18, 3.2, 0.9, [(price, 34, WHITE, True, 1), ("per month", 11, MUTED_DARK, False, 0)])
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.32), Inches(4.35), Inches(3.2), Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = CARD_EDGE; ln.line.fill.background(); ln.shadow.inherit = False
    text(s, x + 0.32, 4.5, 3.3, 1.1, [("\u2022  " + f, 11.5, WHITE, False, 6) for f in feats])
    x += 4.1
text(s, 1.2, 6.05, 10.9, 0.4, [("Core launches with the public rollout; higher tiers add speed, flexibility and multiple accounts as usage deepens.", 12, FAINT, False, 0)], align=PP_ALIGN.CENTER)
text(s, 2.0, 6.45, 9.3, 0.35, [("Subscription model drives predictable recurring revenue.", 11, MUTED_DARK, False, 0)], align=PP_ALIGN.CENTER)
footer(s, 7, light=False)

# ================= SLIDE 08 — COMPETITION =================
s = slide(DARKBG)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("COMPETITION"), 12, LAVENDER, True, 0)])
text(s, 0.62, 1.0, 11.8, 0.9, [("A Capability Gap in the Market", 34, WHITE, True, 0)])
cols = [("", 0.75, 4.35), ("DIME TIME", 5.1, 2.0), ("ROUND-UP APPS", 7.1, 1.85), ("DEBT APPS", 8.95, 1.8), ("BUDGETING APPS", 10.75, 1.85)]
rows = [
    ("Purchase round-ups", "Yes", "Often", "Limited", "Limited"),
    ("Automated debt payments", "Yes", "Rare", "Varies", "No"),
    ("User-controlled allocation", "Yes", "Varies", "Limited", "No"),
    ("Real ACH execution", "Yes", "Varies", "Varies", "No"),
    ("Debt-first positioning", "Yes", "No", "Yes", "No"),
]
top, rh = 2.5, 0.62
# highlight panel behind Dime Time column
hp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(2.05), Inches(2.0), Inches(3.58))
hp.adjustments[0] = 0.08
hp.fill.solid(); hp.fill.fore_color.rgb = LAVENDER; hp.line.fill.background(); hp.shadow.inherit = False
for name, cx, cw in cols[1:]:
    cc = WHITE if cx == 5.1 else LAVENDER
    text(s, cx - 0.1, top - 0.38, cw + 0.2, 0.35, [(name, 10, cc, True, 0)], align=PP_ALIGN.CENTER)
y = top + 0.12
for r in rows:
    text(s, 0.75, y, 4.3, 0.4, [(r[0], 13.5, WHITE, True, 0)])
    for (name, cx, cw), val in zip(cols[1:], r[1:]):
        c = WHITE if cx == 5.1 else FAINT
        text(s, cx, y, cw, 0.4, [(val, 13, c, cx == 5.1, 0)], align=PP_ALIGN.CENTER)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(y + 0.48), Inches(11.85), Pt(0.8))
    ln.fill.solid(); ln.fill.fore_color.rgb = CARD_EDGE; ln.line.fill.background(); ln.shadow.inherit = False
    y += rh
text(s, 1.2, y + 0.25, 10.9, 0.4, [("Dime Time combines automated round-ups, debt-first allocation and real ACH execution in one product.", 13, WHITE, True, 0)], align=PP_ALIGN.CENTER)
text(s, 1.2, y + 0.62, 10.9, 0.35, [("Categories: round-up savings apps, debt-paydown apps, budgeting and financial-wellness apps.", 10, MUTED_DARK, False, 0)], align=PP_ALIGN.CENTER)
footer(s, 8, light=False)

# ================= SLIDE 09 — DEFENSIBILITY =================
s = slide(LAVENDER)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("DEFENSIBILITY"), 12, WHITE, True, 0)])
text(s, 0.62, 1.0, 9.5, 0.9, [("Built to Be Difficult to Replicate", 34, WHITE, True, 0)])
bullets = [
    "Patent-pending dynamic allocation workflow (USPTO provisional filed)",
    "Production ACH infrastructure — real bank-to-bank loop already proven",
    "Compliance and partner-approval track record with regulated providers",
    "Repayment-behavior data advantage as the user base grows",
    "Debt-first brand, user trust and future distribution partnerships",
]
text(s, 0.95, 2.25, 6.3, 4.0, [("\u2022  " + b, 15, WHITE, False, 14) for b in bullets])
pnl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(2.25), Inches(5.15), Inches(4.25))
pnl.fill.background()
pnl.line.color.rgb = WHITE; pnl.line.width = Pt(1.2); pnl.shadow.inherit = False
text(s, 7.8, 2.5, 4.6, 0.4, [(spaced("INFRASTRUCTURE STACK"), 11, WHITE, True, 0)])
stack = [
    ("PLAID", "Bank connections & verification — production access"),
    ("STRIPE", "ACH rails; Financial Connections under review"),
    ("MERCURY", "Business banking & treasury"),
    ("COINBASE", "Crypto execution — planned"),
]
y = 3.0
for name, desc in stack:
    chipsh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.78), Inches(y), Inches(4.6), Inches(0.78))
    chipsh.fill.solid(); chipsh.fill.fore_color.rgb = LAV_PANEL
    chipsh.line.fill.background(); chipsh.shadow.inherit = False
    tf = chipsh.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = name
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(10); r2.font.color.rgb = WHITE
    y += 0.87
text(s, 1.2, 6.58, 10.9, 0.4, [("Proprietary workflow design + production financial infrastructure + partner approvals + a debt-first consumer brand.", 11.5, LAV_SOFT, False, 0)], align=PP_ALIGN.CENTER)
text(s, 7.0, 6.98, 4.3, 0.3, [("Partner names indicate integrations, not endorsements.", 8, LAV_SOFT, False, 0)], align=PP_ALIGN.RIGHT)
footer(s, 9, light=True)

# ================= SLIDE 10 — GO-TO-MARKET =================
s = slide(DARKBG)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("GO-TO-MARKET"), 12, LAVENDER, True, 0)])
text(s, 0.62, 1.0, 11.5, 0.9, [("Focused Launch. Measurable Growth.", 34, WHITE, True, 0)])
# left: who + launch sequence
card(s, 0.72, 2.0, 5.75, 1.7)
text(s, 1.04, 2.25, 5.1, 0.35, [(spaced("WHO WE SERVE FIRST"), 11, LAVENDER, True, 0)])
text(s, 1.04, 2.65, 5.15, 0.95, [("Mobile-first consumers carrying credit-card debt who want to make extra payments but struggle with consistency.", 12.5, WHITE, False, 0)])
card(s, 0.72, 3.9, 5.75, 2.6)
text(s, 1.04, 4.12, 5.1, 0.35, [(spaced("LAUNCH SEQUENCE"), 11, LAVENDER, True, 0)])
steps = [
    "Activate public bank linking",
    "Founder-led launch communications",
    "Acquire a controlled first cohort",
    "Measure activation, retention, ACH volume, conversion",
    "Concentrate spend on channels that prove efficient",
]
text(s, 1.04, 4.5, 5.25, 1.95, [(f"{i+1}.  {st}", 11.5, WHITE, False, 6) for i, st in enumerate(steps)])
# right: channels
card(s, 6.85, 2.0, 5.75, 4.5)
text(s, 7.17, 2.25, 5.1, 0.35, [(spaced("INITIAL CHANNELS"), 11, LAVENDER, True, 0)])
chans = [
    "Founder-led LinkedIn and earned media",
    "Veteran and military communities",
    "Financial-wellness creators",
    "University and young-professional communities",
    "Employer financial-wellness partnerships",
    "Referral incentives once retention is validated",
]
text(s, 7.17, 2.7, 5.2, 3.6, [("\u2022  " + ch, 12.5, WHITE, False, 12) for ch in chans])
text(s, 1.2, 6.7, 10.9, 0.35, [("No paid-spend assumptions — channels are validated before scaling. All growth figures will be reported, not projected.", 10.5, MUTED_DARK, False, 0)], align=PP_ALIGN.CENTER)
footer(s, 10, light=False)

# ================= SLIDE 11 — TRACTION =================
s = slide(DARKBG)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("TRACTION"), 12, LAVENDER, True, 0)])
text(s, 0.62, 1.0, 11.5, 0.9, [("Live App. Real Money Moved.", 34, WHITE, True, 0)])
cards = [
    ("STATUS", ["Live on the Apple", "App Store — v1.0.5", "approved July 2026"]),
    ("MONEY LOOP", ["$1.00 ACH via Stripe;", "$0.99 payout in Mercury —", "proven bank-to-bank"]),
    ("PIPELINE", ["Google Play, bank linking", "& debt import awaiting", "partner approvals"]),
]
x = 0.72
for kick, lines in cards:
    card(s, x, 2.1, 3.86, 2.45)
    text(s, x + 0.3, 2.4, 3.3, 0.35, [(spaced(kick), 11, LAVENDER, True, 0)])
    text(s, x + 0.3, 2.8, 3.35, 1.6, [(ln_, 15, WHITE, True, 3) for ln_ in lines])
    x += 4.1
labels = ["USERS", "LINKED ACCOUNTS", "TRANSACTIONS", "VOLUME", "DEBT REDUCED"]
x = 0.72
for label in labels:
    text(s, x, 5.0, 2.5, 0.35, [(spaced(label), 9.5, LAVENDER, True, 0)], align=PP_ALIGN.CENTER)
    text(s, x, 5.32, 2.5, 0.4, [("Public launch", 13, WHITE, True, 1), ("pending", 13, WHITE, True, 0)], align=PP_ALIGN.CENTER)
    x += 2.42
text(s, 2.0, 6.3, 9.3, 0.4, [(spaced("LIVE ON THE APP STORE  ·  V1.0.5  ·  BUILD 207"), 12, LAVENDER, True, 0)], align=PP_ALIGN.CENTER)
text(s, 2.0, 6.65, 9.3, 0.35, [("No user metrics are reported until public launch — nothing on this slide is estimated.", 10, MUTED_DARK, False, 0)], align=PP_ALIGN.CENTER)
footer(s, 11, light=False)

# ================= SLIDE 12 — TEAM =================
s = slide(DARKBG)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("THE TEAM"), 12, LAVENDER, True, 0)])
text(s, 0.62, 1.0, 11.5, 0.9, [("Focused on Execution", 34, WHITE, True, 0)])
card(s, 2.5, 1.9, 8.3, 3.9)
av = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.95), Inches(2.3), Inches(1.35), Inches(1.35))
av.fill.solid(); av.fill.fore_color.rgb = LAVENDER; av.line.fill.background(); av.shadow.inherit = False
tf = av.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "TC"; r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE
text(s, 4.65, 2.5, 5.5, 0.5, [("Tim Carlisle", 26, WHITE, True, 2)])
text(s, 4.65, 3.15, 5.5, 0.35, [(spaced("FOUNDER & CEO"), 12, LAVENDER, True, 0)])
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.95), Inches(3.95), Inches(7.4), Pt(1))
ln.fill.solid(); ln.fill.fore_color.rgb = CARD_EDGE; ln.line.fill.background(); ln.shadow.inherit = False
team_bullets = [
    "U.S. Navy veteran",
    "MBA — Management Information Systems",
    "Designed, built and shipped Dime Time end-to-end",
    "Integrated Plaid, Stripe ACH and Mercury (Coinbase planned)",
]
text(s, 2.95, 4.15, 7.4, 1.6, [("\u2022  " + b, 12.5, WHITE, False, 8) for b in team_bullets])
text(s, 1.2, 6.1, 10.9, 0.4, [("Solo founder who designed, built and launched a production fintech application with working ACH infrastructure.", 13, FAINT, False, 0)], align=PP_ALIGN.CENTER)
footer(s, 12, light=False)

# ================= SLIDE 13 — THE RAISE =================
s = slide(LAVENDER)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("THE RAISE"), 12, WHITE, True, 0)])
text(s, 0.62, 1.0, 11.8, 0.9, [("$150K to Launch and Prove Traction", 34, WHITE, True, 0)])
# left: stats
stats = [("RAISING", "$150,000"), ("INSTRUMENT", "Post-money SAFE"), ("EST. RUNWAY", "~12 months")]
y = 2.1
for k, v in stats:
    c = card(s, 0.72, y, 3.4, 1.28, dark=False)
    text(s, 1.0, y + 0.18, 2.9, 0.3, [(spaced(k), 9.5, WHITE, True, 0)])
    text(s, 1.0, y + 0.5, 2.95, 0.55, [(v, 21 if len(v) < 12 else 17, WHITE, True, 0)])
    y += 1.48
# middle: use of funds
text(s, 4.55, 2.15, 4.0, 0.35, [(spaced("USE OF FUNDS"), 11, WHITE, True, 0)])
funds = [(35, "Marketing & user acquisition"), (30, "Founder runway"),
         (25, "Engineering & contractor support"), (10, "Compliance, infra & operations")]
y = 2.65
for pct, lbl in funds:
    text(s, 4.55, y, 0.95, 0.4, [(f"{pct}%", 17, WHITE, True, 0)])
    text(s, 5.5, y + 0.05, 3.0, 0.4, [(lbl, 11, WHITE, False, 0)])
    barbg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(y + 0.52), Inches(3.85), Inches(0.16))
    barbg.adjustments[0] = 0.5
    barbg.fill.solid(); barbg.fill.fore_color.rgb = LAV_PANEL; barbg.line.fill.background(); barbg.shadow.inherit = False
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(y + 0.52), Inches(3.85 * pct / 35.0), Inches(0.16))
    bar.adjustments[0] = 0.5
    bar.fill.solid(); bar.fill.fore_color.rgb = WHITE; bar.line.fill.background(); bar.shadow.inherit = False
    y += 0.98
# right: milestones
c = card(s, 8.85, 2.1, 3.78, 4.35, dark=False)
text(s, 9.15, 2.35, 3.2, 0.35, [(spaced("MILESTONES"), 11, WHITE, True, 0)])
miles = [
    "Launch publicly on iOS and Android",
    "Prove repeatable customer acquisition",
    "Reach an initial paying-subscriber target",
    "Demonstrate recurring ACH payment volume",
    "Set up a larger seed round with measurable traction",
]
text(s, 9.15, 2.8, 3.25, 3.4, [("\u2022  " + m, 11.5, WHITE, False, 11) for m in miles])
text(s, 1.2, 6.7, 10.9, 0.35, [("SAFE terms available on request. Allocation percentages total 100%.", 10.5, LAV_SOFT, False, 0)], align=PP_ALIGN.CENTER)
footer(s, 13, light=True)

# ================= SLIDE 14 — COMPANY INFORMATION =================
s = slide(DARKBG)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("APPENDIX"), 12, LAVENDER, True, 0)])
text(s, 0.62, 1.0, 11.5, 0.9, [("Company Information", 34, WHITE, True, 0)])
info = [("COMPANY", "Dime Time Technologies LLC"), ("FOUNDER", "Tim Carlisle"),
        ("WEBSITE", "dime-time.com"), ("EMAIL", "tim@dime-time.com")]
y = 2.4
for k, v in info:
    text(s, 0.72, y, 3.5, 0.4, [(spaced(k), 12, LAVENDER, True, 0)])
    text(s, 5.0, y - 0.05, 7.6, 0.45, [(v, 17, WHITE, True, 0)], align=PP_ALIGN.RIGHT)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(y + 0.5), Inches(11.9), Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = CARD_EDGE; ln.line.fill.background(); ln.shadow.inherit = False
    y += 0.85
footer(s, 14, light=False)

# ================= SAVE =================
# NOTE: order matters — the main 14-slide deck MUST be saved BEFORE the
# with-ip block below, which appends 3 IP slides to the same `prs` object.
main_out = os.path.join(HERE, "dime-time-pitch-deck.pptx")
prs.save(main_out)
print("saved", main_out, "-", len(prs.slides._sldIdLst), "slides")

# with-ip variant: append the 3 IP appendix PNG slides
for ip in ["appendix-ip-1-patent.png", "appendix-ip-2-architecture.png", "appendix-ip-3-ip-claims.png"]:
    p = os.path.join(HERE, ip)
    sl = prs.slides.add_slide(BLANK)
    sl.shapes.add_picture(p, 0, 0, width=W, height=H)
ip_out = os.path.join(HERE, "dime-time-pitch-deck-with-ip.pptx")
prs.save(ip_out)
print("saved", ip_out)
