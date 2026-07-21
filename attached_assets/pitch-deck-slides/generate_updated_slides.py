"""Regenerate pitch deck slides 09 (defensibility) + 10 (traction) with July 2026 facts.

Renders a 2-slide pptx; render pipeline (soffice -> pdftoppm) replaces the
slide-09/slide-10 PNGs, then the full deck pptx/pdf is rebuilt from PNG stack.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

LAVENDER = RGBColor(0x91, 0x8E, 0xF4)
LAV_PANEL = RGBColor(0xA4, 0xA1, 0xF6)
DARKBG = RGBColor(0x21, 0x1D, 0x3D)
CARD = RGBColor(0x2B, 0x27, 0x4E)
CARD_EDGE = RGBColor(0x4A, 0x46, 0x8F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FAINT = RGBColor(0xC9, 0xC7, 0xF9)
DEEP = RGBColor(0x4A, 0x46, 0x8F)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def slide(bg):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    for i, (t, size, color, bold, sp) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp)
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tf


def spaced(t):
    return " ".join(list(t.replace(" ", "  ")))


def footer(s, n, light):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(7.06), Inches(0.42), Pt(1.6))
    c = WHITE if not light else WHITE
    ln.fill.solid(); ln.fill.fore_color.rgb = c
    ln.line.fill.background(); ln.shadow.inherit = False
    text(s, 1.08, 6.9, 3.0, 0.35, [(spaced("DIME TIME"), 11, c, True, 0)])
    text(s, 11.8, 6.95, 1.1, 0.35, [(f"{n} / 13", 11, FAINT if not light else RGBColor(0xE3, 0xE2, 0xFC), False, 0)], align=PP_ALIGN.RIGHT)


SHOTS = os.path.join(os.path.dirname(__file__), "..", "Dime-Time-App-Store-Screenshots-FINAL", "iPhone-6.9")


def phone(s, x, y, h, png):
    """Dark rounded device frame + screenshot inset. Aspect 1290x2796."""
    w = h * 1290 / 2796
    pad = 0.07
    fr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - pad), Inches(y - pad),
                            Inches(w + 2 * pad), Inches(h + 2 * pad))
    fr.fill.solid(); fr.fill.fore_color.rgb = RGBColor(0x14, 0x12, 0x28)
    fr.line.fill.background(); fr.shadow.inherit = False
    fr.adjustments[0] = 0.09
    s.shapes.add_picture(png, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    return w


# ---------- SLIDE 03 — SOLUTION ----------
s = slide(LAVENDER)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("THE SOLUTION"), 12, WHITE, True, 0)])
text(s, 0.62, 1.0, 9.0, 0.9, [("Automated Debt Reduction", 34, WHITE, True, 0)])
sol_bullets = [
    "Every purchase rounds up automatically",
    "Spare change is applied to your financial goals",
    "No behavior change required",
    "Works passively in the background",
]
text(s, 0.95, 2.3, 7.6, 3.4, [("\u2022  " + b, 17, WHITE, False, 20) for b in sol_bullets])
phone(s, 9.85, 1.35, 5.35, os.path.join(SHOTS, "01-dashboard-iphone.png"))
footer(s, "03", light=True)

# ---------- SLIDE 05 — PRODUCT ----------
s = slide(DARKBG)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("THE PRODUCT"), 12, LAVENDER, True, 0)])
text(s, 0.62, 1.0, 11.5, 0.9, [("Live Product — Real Financial Movement", 32, WHITE, True, 0)])
shots = [
    ("01-dashboard-iphone.png", "Live dashboard"),
    ("02-debts-iphone.png", "Debt tracking"),
    ("03-roundups-iphone.png", "Round-ups"),
    ("04-insights-iphone.png", "Insights"),
]
ph_h = 3.85
ph_w = ph_h * 1290 / 2796
gap = (13.333 - 2 * 1.15 - 4 * ph_w) / 3
x = 1.15
for png, label in shots:
    phone(s, x, 2.0, ph_h, os.path.join(SHOTS, png))
    text(s, x - 0.45, 6.0, ph_w + 0.9, 0.4, [(label, 14, WHITE, True, 0)], align=PP_ALIGN.CENTER)
    x += ph_w + gap
text(s, 2.0, 6.55, 9.3, 0.35, [("As shipped in v1.0.5 — live on the App Store, July 2026", 12, FAINT, False, 0)], align=PP_ALIGN.CENTER)
footer(s, "05", light=False)

# ---------- SLIDE 09 — DEFENSIBILITY ----------
s = slide(LAVENDER)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("DEFENSIBILITY"), 12, WHITE, True, 0)])
text(s, 0.62, 1.0, 9.5, 0.9, [("Built to Be Difficult to Replicate", 34, WHITE, True, 0)])
bullets = [
    "Patent-pending allocation engine",
    "USPTO provisional patent filed — patent pending",
    "Real-time financial routing system",
    "Live on the Apple App Store (v1.0.5)",
    "Real ACH money movement proven in production",
]
tf = text(s, 0.95, 2.25, 6.3, 4.0, [("\u2022  " + b, 16, WHITE, False, 14) for b in bullets])
pnl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(2.25), Inches(5.15), Inches(4.25))
pnl.fill.background()
pnl.line.color.rgb = WHITE; pnl.line.width = Pt(1.2); pnl.shadow.inherit = False
text(s, 7.8, 2.5, 4.6, 0.4, [(spaced("INFRASTRUCTURE STACK"), 11, WHITE, True, 0)])
stack = [
    ("PLAID", "Bank connections & verification"),
    ("STRIPE", "ACH rails & Financial Connections"),
    ("MERCURY", "Business banking & treasury"),
    ("COINBASE", "Crypto execution & custody"),
]
y = 3.0
for name, desc in stack:
    chipsh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.78), Inches(y), Inches(4.6), Inches(0.78))
    chipsh.fill.solid(); chipsh.fill.fore_color.rgb = LAV_PANEL
    chipsh.line.fill.background(); chipsh.shadow.inherit = False
    tf = chipsh.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = name
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(10.5); r2.font.color.rgb = WHITE
    y += 0.87
text(s, 2.0, 6.9, 9.3, 0.4, [("Replicating this stack requires 6 – 12 months of compliance and engineering.", 12, RGBColor(0xE9, 0xE8, 0xFD), False, 0)], align=PP_ALIGN.CENTER)
footer(s, "09", light=True)

# ---------- SLIDE 10 — TRACTION ----------
s = slide(DARKBG)
text(s, 0.68, 0.62, 8, 0.4, [(spaced("TRACTION"), 12, LAVENDER, True, 0)])
text(s, 0.62, 1.0, 11.5, 0.9, [("Live App. Real Money Moved.", 34, WHITE, True, 0)])
cards = [
    ("STATUS", ["Live on the Apple", "App Store — v1.0.5", "approved July 2026"]),
    ("MONEY LOOP", ["$1.00 ACH via Stripe;", "payout in Mercury —", "proven bank-to-bank"]),
    ("PIPELINE", ["Google Play, bank linking", "& debt import awaiting", "partner approvals"]),
]
x = 0.72
for kick, lines in cards:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.1), Inches(3.86), Inches(2.55))
    card.fill.solid(); card.fill.fore_color.rgb = CARD
    card.line.color.rgb = CARD_EDGE; card.line.width = Pt(1); card.shadow.inherit = False
    text(s, x + 0.3, 2.45, 3.3, 0.35, [(spaced(kick), 11, LAVENDER, True, 0)])
    text(s, x + 0.3, 2.85, 3.3, 1.6, [(ln, 15.5, WHITE, True, 3) for ln in lines])
    x += 4.1
metrics = [("USERS", "\u2014"), ("LINKED ACCOUNTS", "\u2014"), ("TRANSACTIONS", "\u2014"),
           ("TOTAL VOLUME", "$\u2014"), ("DEBT REDUCED", "$\u2014")]
x = 0.72
for label, val in metrics:
    text(s, x, 5.05, 2.5, 0.35, [(spaced(label), 9.5, LAVENDER, True, 0)], align=PP_ALIGN.CENTER)
    text(s, x, 5.35, 2.5, 0.5, [(val, 22, WHITE, True, 0)], align=PP_ALIGN.CENTER)
    x += 2.42
text(s, 2.0, 6.25, 9.3, 0.4, [(spaced("LIVE ON THE APP STORE  ·  V1.0.5  ·  BUILD 207"), 12, LAVENDER, True, 0)], align=PP_ALIGN.CENTER)
text(s, 2.0, 6.6, 9.3, 0.35, [("Replace dashes with live numbers before presenting.", 10, RGBColor(0x8B, 0x88, 0xB0), False, 0)], align=PP_ALIGN.CENTER)
footer(s, "10", light=False)

out = os.path.join(os.path.dirname(__file__), "_updated_slides.pptx")
prs.save(out)
print("saved", out)
